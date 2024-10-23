from pptx import Presentation

# Create a new presentation
presentation = Presentation()

# Slide 1: Title Slide
slide_1 = presentation.slides.add_slide(presentation.slide_layouts[0])
title_1 = slide_1.shapes.title
subtitle_1 = slide_1.placeholders[1]
title_1.text = "Design Issues of Subprograms, Arithmetic Expressions, and Overloaded Operators"
subtitle_1.text = "An In-Depth Exploration with Examples"

# Slide 2: Introduction to Subprograms
slide_2 = presentation.slides.add_slide(presentation.slide_layouts[1])
title_2 = slide_2.shapes.title
content_2 = slide_2.placeholders[1]
title_2.text = "Introduction to Subprograms"
content_2.text = (
    "• Subprograms (functions, methods) are self-contained code blocks.\n"
    "• Importance:\n"
    "  - Modularity: Breaks down complex problems.\n"
    "  - Reusability: Same code used multiple times.\n"
    "  - Maintenance: Easier to debug.\n"
    "• Example:\n"
    "  def square(num): return num * num"
)

# Slide 3: Parameter Passing Techniques
slide_3 = presentation.slides.add_slide(presentation.slide_layouts[1])
title_3 = slide_3.shapes.title
content_3 = slide_3.placeholders[1]
title_3.text = "Parameter Passing Techniques"
content_3.text = (
    "• How arguments are passed to subprograms affects their behavior.\n"
    "• Types of Parameter Passing:\n"
    "  - Pass-by-Value: Copy passed; original not affected.\n"
    "    Example (C):\n"
    "      void increment(int a) { a++; }\n"
    "  - Pass-by-Reference: Reference passed; original affected.\n"
    "    Example (C++):\n"
    "      void increment(int &a) { a++; }"
)

# Slide 4: Return Values in Subprograms
slide_4 = presentation.slides.add_slide(presentation.slide_layouts[1])
title_4 = slide_4.shapes.title
content_4 = slide_4.placeholders[1]
title_4.text = "Return Values in Subprograms"
content_4.text = (
    "• Subprograms can return values influencing result handling.\n"
    "• Single Return Value:\n"
    "  Example (Python): def add(a, b): return a + b\n"
    "• Multiple Return Values:\n"
    "  Example (Python): def get_info(): return 'Alice', 30\n"
    "• Complex Objects:\n"
    "  Example (Python): def get_dict(): return {'name': 'Alice', 'age': 30}"
)

# Slide 5: Subprogram Overloading
slide_5 = presentation.slides.add_slide(presentation.slide_layouts[1])
title_5 = slide_5.shapes.title
content_5 = slide_5.placeholders[1]
title_5.text = "Subprogram Overloading"
content_5.text = (
    "• Overloading allows multiple subprograms with the same name.\n"
    "• Benefits: Improves readability and usability.\n"
    "• Example (C++):\n"
    "  void display(int a);\n"
    "  void display(double b);\n"
    "  void display(string c);"
)

# Slide 6: Understanding Recursion
slide_6 = presentation.slides.add_slide(presentation.slide_layouts[1])
title_6 = slide_6.shapes.title
content_6 = slide_6.placeholders[1]
title_6.text = "Understanding Recursion"
content_6.text = (
    "• Recursion occurs when a subprogram calls itself.\n"
    "• Components:\n"
    "  - Base Case: Condition to stop recursion.\n"
    "  - Recursive Case: Calls itself with modified arguments.\n"
    "• Example (Python):\n"
    "  def factorial(n):\n"
    "      if n == 0:\n"
    "          return 1\n"
    "      else:\n"
    "          return n * factorial(n - 1)"
)

# Slide 7: Scope and Lifetime of Variables
slide_7 = presentation.slides.add_slide(presentation.slide_layouts[1])
title_7 = slide_7.shapes.title
content_7 = slide_7.placeholders[1]
title_7.text = "Scope and Lifetime of Variables"
content_7.text = (
    "• Scope: Visibility of variables in a program.\n"
    "  - Local Scope: Variables within a function.\n"
    "    Example (Python):\n"
    "      def func():\n"
    "          x = 10\n"
    "  - Global Scope: Accessible throughout the program.\n"
    "    Example (Python):\n"
    "      x = 10\n"
    "      def func():\n"
    "          print(x)"
)

# Slide 8: Handling Side Effects
slide_8 = presentation.slides.add_slide(presentation.slide_layouts[1])
title_8 = slide_8.shapes.title
content_8 = slide_8.placeholders[1]
title_8.text = "Handling Side Effects"
content_8.text = (
    "• Side effects occur when a subprogram alters some state outside its scope.\n"
    "• Impact: Can lead to unexpected behavior.\n"
    "• Example (C):\n"
    "  int global_var = 10;\n"
    "  void modify() { global_var++; }\n"
    "• Best Practices: Minimize side effects by using local variables."
)

# Slide 9: Exception Handling Mechanisms
slide_9 = presentation.slides.add_slide(presentation.slide_layouts[1])
title_9 = slide_9.shapes.title
content_9 = slide_9.placeholders[1]
title_9.text = "Exception Handling Mechanisms"
content_9.text = (
    "• Exception handling allows a program to respond to runtime errors.\n"
    "• Importance: Prevents crashes and alternative execution paths.\n"
    "• Example (Python):\n"
    "  try:\n"
    "      result = 10 / 0\n"
    "  except ZeroDivisionError:\n"
    "      print('Cannot divide by zero!')"
)

# Slide 10: Arithmetic Expressions and Operator Overloading
slide_10 = presentation.slides.add_slide(presentation.slide_layouts[1])
title_10 = slide_10.shapes.title
content_10 = slide_10.placeholders[1]
title_10.text = "Arithmetic Expressions and Operator Overloading"
content_10.text = (
    "• Arithmetic expressions use operators (+, -, *, /) for calculations.\n"
    "• Operator Overloading: Customizing how operators work with user-defined types.\n"
    "• Example (Python):\n"
    "  class Vector:\n"
    "      def __init__(self, x, y):\n"
    "          self.x = x\n"
    "          self.y = y\n"
    "      def __add__(self, other):\n"
    "          return Vector(self.x + other.x, self.y + other.y)\n"
    "  v1 = Vector(1, 2)\n"
    "  v2 = Vector(3, 4)\n"
    "  result = v1 + v2"
)

# Slide 11: Conclusion
slide_11 = presentation.slides.add_slide(presentation.slide_layouts[1])
title_11 = slide_11.shapes.title
content_11 = slide_11.placeholders[1]
title_11.text = "Conclusion"
content_11.text = (
    "• Design of subprograms, arithmetic expressions, and operator overloading are critical for programming.\n"
    "• Key Takeaways:\n"
    "  - Proper parameter passing enhances clarity.\n"
    "  - Understanding recursion and scope improves reliability.\n"
    "  - Minimizing side effects contributes to robust code."
)

# Save the presentation to a file
pptx_file_path = '/mnt/data/Design_Issues_Subprograms_Arithemetic_Expressions_Overloaded_Operators.pptx'
presentation.save(pptx_file_path)

pptx_file_path
